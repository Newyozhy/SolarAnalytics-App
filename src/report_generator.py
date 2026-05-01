import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os
import tempfile

def generate_excel(dataframes, output_path):
    """
    Genera un archivo Excel multi-hoja a partir del diccionario de DataFrames.
    """
    if not dataframes:
        raise ValueError("No hay datos para exportar.")
        
    with pd.ExcelWriter(output_path, engine='xlsxwriter') as writer:
        for name, df in dataframes.items():
            if df is not None and not df.empty:
                # Truncar el nombre de la hoja a 31 caracteres máximo
                sheet_name = str(name)[:31]
                df.to_excel(writer, sheet_name=sheet_name, index=False)

def generate_pptx(template_path, output_path, project_name, figures=None, data_summary=None):
    """
    Toma la plantilla PPTX, inserta el nombre del proyecto, datos básicos
    y reemplaza o inserta las gráficas generadas (figures).
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"No se encontró la plantilla en {template_path}")
        
    prs = Presentation(template_path)
    
    # Reemplazar texto simple en todas las diapositivas
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{PROJECT_NAME}" in run.text:
                        run.text = run.text.replace("{PROJECT_NAME}", project_name)
                    
                    if data_summary:
                        if "{TOTAL_GEN}" in run.text and 'total_gen' in data_summary:
                            run.text = run.text.replace("{TOTAL_GEN}", f"{data_summary['total_gen']:,.2f}")
                        if "{AVG_GEN}" in run.text and 'avg_gen' in data_summary:
                            run.text = run.text.replace("{AVG_GEN}", f"{data_summary['avg_gen']:,.2f}")

    # Si hay gráficas, las insertamos en una diapositiva nueva al final
    if figures:
        blank_slide_layout = prs.slide_layouts[6] # Suele ser una en blanco
        for name, fig in figures.items():
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Título de la gráfica
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = f"Gráfica: {name}"
            
            # Exportar gráfica temporalmente
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                temp_file = tmp.name
                
            try:
                # Requiere kaleido instalado
                fig.write_image(temp_file)
                # Insertar imagen en la diapositiva
                slide.shapes.add_picture(temp_file, Inches(1), Inches(1.5), height=Inches(5))
            finally:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    
    prs.save(output_path)
    return output_path
